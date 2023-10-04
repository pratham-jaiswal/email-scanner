import tkinter as tk
from tkinter import ttk, messagebox
import imaplib
import email
from datetime import datetime
import os
import subprocess
from PyPDF2 import PdfReader
import docx
from pptx import Presentation
import pandas as pd
import json
import pytz
import threading

def app():

    def load_search_terms():
        try:
            with open("search_terms.json", "r") as file:
                data = json.load(file)
                return data["search_terms"]
        except FileNotFoundError:
            search_terms = {"search_terms": ["Your Name", "youremail@example.com", "Placement", "Internship", "whatever you like"]}
            with open('search_terms.json', 'w') as file:
                json.dump(search_terms, file, indent=4)
            return search_terms
    
    def edit_search_terms():
        try:
            subprocess.Popen(["notepad.exe", "search_terms.json"])
        except Exception as e:
            messagebox.showerror("Error", f"An error occurred: {str(e)}")

    def create_menu():
        menu = tk.Menu(window)
        window.config(menu=menu)

        edit_menu = tk.Menu(menu, tearoff=0)
        menu.add_cascade(label="Configure", menu=edit_menu)
        edit_menu.add_command(label="Edit Search Terms", command=edit_search_terms)

    def search_in_pdf(file_path, search_terms):
        try:
            pdf = PdfReader(file_path)
            for page in pdf.pages:
                page_text = page.extract_text()
                for term in search_terms:
                    if term in page_text:
                        return True
            return False
        except Exception:
            return False

    def search_in_docx(file_path, search_terms):
        try:
            doc = docx.Document(file_path)
            for paragraph in doc.paragraphs:
                for term in search_terms:
                    if term in paragraph.text:
                        return True
            return False
        except Exception:
            return False

    def search_in_pptx(file_path, search_terms):
        try:
            presentation = Presentation(file_path)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        for term in search_terms:
                            if term in shape.text:
                                return True
            return False
        except Exception:
            return False

    def search_in_csv_xlsx(file_path, search_terms):
        try:
            df = pd.read_excel(file_path)
            for column in df.columns:
                for term in search_terms:
                    if df[column].astype(str).str.contains(term).any():
                        return True
            return False
        except Exception:
            return False

    def download_attachments(email_message, save_directory):
        for part in email_message.walk():
            if part.get_content_maintype() == "application":
                filename = part.get_filename()
                if filename:
                    file_path = os.path.join(save_directory, filename)
                    with open(file_path, 'wb') as file:
                        file.write(part.get_payload(decode=True))
                    yield file_path

    def update_scans(df, subject, fromAddr, dt, scans):
        fromAddr = list(fromAddr)
        
        if not ((df['Subject'] == subject) & (df['Datetime'] == dt)).any():
            data = pd.DataFrame({'Subject': subject, 'From': fromAddr, 'Datetime': dt, 'Checked': "No"})
            df = pd.concat([df, data], ignore_index=True)
            df.to_csv(scans, index=False)

    def scan_emails(imap, msgnums, df, scans, save_directory):
        search_terms = load_search_terms()
        for msgnum in msgnums[0].split():
            found = False

            _, data = imap.fetch(msgnum, "(RFC822)")
            message = email.message_from_bytes(data[0][1])

            body = ''
            for part in message.walk():
                if part.get_content_type() == "text/plain":
                    body += part.as_string()
            
            for term in search_terms:
                if term.lower() in body.lower():
                    update_scans(df, message.get('Subject'), {message.get('From')}, message.get('Date'), scans)

            for term in search_terms:
                for attachment_path in download_attachments(message, save_directory):
                    found = search_in_pdf(attachment_path, search_terms) or \
                        search_in_docx(attachment_path, search_terms) or \
                        search_in_pptx(attachment_path, search_terms) or \
                        search_in_csv_xlsx(attachment_path, search_terms)

                    if found:
                        break

                if found:
                    update_scans(df, message.get('Subject'), {message.get('From')}, message.get('Date'), scans)

    def delete_attachments(save_directory):
        for attachment_path in os.listdir(save_directory):
            os.remove(os.path.join(save_directory, attachment_path))

        os.rmdir(save_directory)

    def change_status(results, scans, new_status):
        selected_items = results.selection()
        if not selected_items:
            messagebox.showwarning("Warning", "Please select a row to update.")
            return

        df = pd.read_csv(scans)
        for item in selected_items:
            subject = results.item(item, "values")[0]
            dt = results.item(item, "values")[2]
            df.loc[(df['Subject'] == subject) & (df['Datetime'] == dt), 'Checked'] = 'No' if df['Checked'].any() == 'Yes' else 'Yes'

        df.to_csv(scans, index=False)
        refresh(results, scans)

    def delete_entry(results, scans):
        selected_items = results.selection()
        if not selected_items:
            messagebox.showwarning("Warning", "Please select a row to delete.")
            return

        df = pd.read_csv(scans)
        for item in selected_items:
            subject = results.item(item, "values")[0]
            dt = results.item(item, "values")[2]
            df = df[~((df['Subject'] == subject) & (df['Datetime'] == dt))]

        df.to_csv(scans, index=False)
        refresh(results, scans)

    def refresh(results, scans):
        df = pd.read_csv(scans)
        results.delete(*results.get_children())  # Clear existing items
        for _, row in df.iterrows():
            results.insert("", "end", values=(row["Subject"], row["From"], row["Datetime"], row["Checked"]))

    def click_start():
        imap_server = "imap.gmail.com"
        email_address = email_entry.get()
        password = password_entry.get()

        imap = imaplib.IMAP4_SSL(imap_server)

        try:
            imap.login(email_address, password)
            imap.select("INBOX")
            
            user_timezone = pytz.timezone('Asia/Kolkata')
            current_datetime = datetime.now().astimezone(user_timezone)

            since_date = (current_datetime).strftime("%d-%b-%Y")
            _, msgnums = imap.search(None, f'SINCE {since_date}')

            save_directory = "attachments"
            os.makedirs(save_directory, exist_ok=True)
            
            scan_emails(imap, msgnums, df, scans, save_directory)
            delete_attachments(save_directory)

            imap.close()

            refresh(result_data, scans)
            messagebox.showinfo("Success", "Scanning completed.")
        except Exception as e:
            messagebox.showerror("Error", f"An error occurred: {str(e)}")

    def start_thread():
        threading.Thread(target=click_start).start()

    scans = 'scans.csv'
    if not os.path.exists(scans):
        initial_data = {'Subject': [], 'From': [], 'Datetime': [], 'Checked': []}
        df = pd.DataFrame(initial_data)
        df.to_csv(scans, index=False)

    df = pd.read_csv(scans)

    window = tk.Tk()
    window.title("Email Scanner")
    
    search_terms = load_search_terms()
    create_menu()

    style = ttk.Style()
    style.configure("TButton", padding=(10, 5), font=("Helvetica", 12))
    style.configure("TLabel", font=("Helvetica", 12))

    email_label = ttk.Label(window, text="Email Address:")
    email_entry = ttk.Entry(window, width=50, font=("Helvetica", 10))

    password_label = ttk.Label(window, text="Password:")
    password_entry = ttk.Entry(window, show="*", width=50, font=("Helvetica", 10))

    start_button = ttk.Button(window, text="Start Scan", command=lambda: start_thread())

    result_data = ttk.Treeview(window, columns=("Subject", "From", "Datetime", "Checked"), show="headings")
    result_data.heading("Subject", text="Subject")
    result_data.heading("From", text="From")
    result_data.heading("Datetime", text="Datetime")
    result_data.heading("Checked", text="Checked")

    change_status_button = ttk.Button(window, text="Change Status", command=lambda: change_status(result_data, scans, "Yes"))
    delete_button = ttk.Button(window, text="Delete", command=lambda: delete_entry(result_data, scans))
    update_button = ttk.Button(window, text="Refresh", command=lambda: refresh(result_data, scans))

    window.grid_columnconfigure(0, weight=1)
    window.grid_columnconfigure(1, weight=1)
    window.grid_rowconfigure(3, weight=1)

    email_label.grid(row=0, column=0, padx=10, pady=5, sticky="w")
    email_entry.grid(row=0, column=1, padx=10, pady=5)

    password_label.grid(row=1, column=0, padx=10, pady=5, sticky="w")
    password_entry.grid(row=1, column=1, padx=10, pady=5)

    start_button.grid(row=2, column=0, columnspan=2, pady=10)

    result_data.grid(row=3, column=0, columnspan=2, padx=10, pady=5, sticky="nsew")

    change_status_button.grid(row=4, column=0, padx=10, pady=5)
    delete_button.grid(row=4, column=1, padx=10, pady=5)
    update_button.grid(row=5, column=0, columnspan=2, pady=10)

    result_data.grid(row=3, column=0, columnspan=2, padx=10, pady=5, sticky="nsew")

    vsb = ttk.Scrollbar(window, orient="vertical", command=result_data.yview)
    result_data.configure(yscrollcommand=vsb.set)
    vsb.grid(row=3, column=2, sticky="ns")

    refresh(result_data, scans)

    window.resizable(True, True)
    window.mainloop()

if __name__ == "__main__":
    app()